import json
import os
from io import BytesIO
from pathlib import Path

import pandas as pd
import pdfplumber
from langchain_core.documents import Document
from pptx import Presentation

SUPPORTED_FILE_TYPES = {
    ".pdf": "PDF document",
    ".pptx": "PowerPoint presentation",
    ".xlsx": "Excel workbook",
    ".xlsm": "Excel macro workbook",
    ".csv": "CSV dataset",
    ".tsv": "TSV dataset",
    ".txt": "Text document",
    ".md": "Markdown document",
    ".json": "JSON document",
    # Image formats (OCR)
    ".jpg": "JPEG image",
    ".jpeg": "JPEG image",
    ".png": "PNG image",
    ".tiff": "TIFF image",
    ".tif": "TIFF image",
    ".bmp": "BMP image",
    ".webp": "WebP image",
}

TABULAR_EXTENSIONS = {".xlsx", ".xlsm", ".csv", ".tsv"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp"}


def _to_bytes(uploaded_file):
    if hasattr(uploaded_file, "getvalue"):
        return uploaded_file.getvalue()

    if hasattr(uploaded_file, "read"):
        if hasattr(uploaded_file, "seek"):
            uploaded_file.seek(0)
        return uploaded_file.read()

    raise TypeError("Unsupported uploaded file object.")


def _clean_cell(value):
    if pd.isna(value):
        return ""
    return str(value).replace("\n", " ").strip()


def _normalize_dataframe(dataframe):
    normalized_df = dataframe.copy()
    normalized_df.columns = [str(column).strip() for column in normalized_df.columns]
    normalized_df = normalized_df.fillna("")
    return normalized_df


def _read_excel_frames(file_obj):
    workbook_bytes = BytesIO(_to_bytes(file_obj))
    excel_file = pd.ExcelFile(workbook_bytes, engine="openpyxl")
    frames = []
    for sheet_name in excel_file.sheet_names:
        dataframe = excel_file.parse(sheet_name=sheet_name)
        frames.append((sheet_name, _normalize_dataframe(dataframe)))
    return frames


def _read_delimited_frame(file_obj, extension):
    delimiter = "\t" if extension == ".tsv" else ","
    dataframe = pd.read_csv(BytesIO(_to_bytes(file_obj)), sep=delimiter)
    return _normalize_dataframe(dataframe)


def _tabular_chunk_documents(dataframe, source_name, file_type, sheet_name="Data", chunk_size=20):
    documents = []
    normalized_df = _normalize_dataframe(dataframe)

    column_list = [column for column in normalized_df.columns if column]
    row_count = len(normalized_df.index)

    summary = (
        f"Source: {source_name}\n"
        f"Sheet: {sheet_name}\n"
        f"File type: {file_type}\n"
        f"Columns: {', '.join(column_list) if column_list else 'No columns detected'}\n"
        f"Row count: {row_count}"
    )
    documents.append(
        Document(
            page_content=summary,
            metadata={
                "source": source_name,
                "page": f"Sheet {sheet_name}",
                "sheet": sheet_name,
                "file_type": file_type,
                "content_type": "table_summary",
            },
        )
    )

    for start_index in range(0, row_count, chunk_size):
        end_index = min(start_index + chunk_size, row_count)
        chunk_df = normalized_df.iloc[start_index:end_index]
        row_lines = []
        for row_number, (_, row) in enumerate(chunk_df.iterrows(), start=start_index + 1):
            cell_pairs = []
            for column in column_list:
                cleaned = _clean_cell(row[column])
                if cleaned:
                    cell_pairs.append(f"{column}: {cleaned}")
            if cell_pairs:
                row_lines.append(f"Row {row_number}: " + " | ".join(cell_pairs))

        if not row_lines:
            continue

        content = (
            f"Source: {source_name}\n"
            f"Sheet: {sheet_name}\n"
            f"Columns: {', '.join(column_list) if column_list else 'No columns detected'}\n"
            f"Rows: {start_index + 1}-{end_index}\n"
            + "\n".join(row_lines)
        )
        documents.append(
            Document(
                page_content=content,
                metadata={
                    "source": source_name,
                    "page": f"Sheet {sheet_name} Rows {start_index + 1}-{end_index}",
                    "sheet": sheet_name,
                    "row_start": start_index + 1,
                    "row_end": end_index,
                    "file_type": file_type,
                    "content_type": "table_rows",
                },
            )
        )

    return documents


def load_pdf(file_obj, source_name):
    documents = []
    with pdfplumber.open(BytesIO(_to_bytes(file_obj))) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            if text:
                documents.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": source_name,
                            "page": page_num,
                            "file_type": "pdf",
                            "content_type": "pdf_page",
                        },
                    )
                )
    return documents


def load_pptx(file_obj, source_name):
    documents = []
    presentation = Presentation(BytesIO(_to_bytes(file_obj)))

    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                cleaned = shape.text.strip()
                if cleaned:
                    text_parts.append(cleaned)

        if not text_parts:
            continue

        slide_text = "\n".join(text_parts)
        documents.append(
            Document(
                page_content=slide_text,
                metadata={
                    "source": source_name,
                    "page": f"Slide {slide_index}",
                    "slide": slide_index,
                    "file_type": "pptx",
                    "content_type": "slide",
                },
            )
        )

    return documents


def load_excel(file_obj, source_name, extension):
    documents = []
    for sheet_name, dataframe in _read_excel_frames(file_obj):
        documents.extend(_tabular_chunk_documents(dataframe, source_name, extension.lstrip("."), sheet_name))

    return documents


def load_delimited_text(file_obj, source_name, extension):
    dataframe = _read_delimited_frame(file_obj, extension)
    return _tabular_chunk_documents(dataframe, source_name, extension.lstrip("."), "Data")


def load_text_document(file_obj, source_name, extension):
    text = _to_bytes(file_obj).decode("utf-8", errors="ignore")
    if not text.strip():
        return []

    return [
        Document(
            page_content=text,
            metadata={
                "source": source_name,
                "page": extension.lstrip(".").upper(),
                "file_type": extension.lstrip("."),
                "content_type": "text_document",
            },
        )
    ]


def load_image_ocr(file_obj, source_name):
    """Extract text from an image via OCR.

    Tries pytesseract first (requires system Tesseract install), then falls
    back to Gemini Vision if GOOGLE_API_KEY is set. Raises ValueError when
    neither engine is available.
    """
    image_bytes = _to_bytes(file_obj)
    extension = Path(source_name).suffix.lower()

    # ── pytesseract ───────────────────────────────────────────────────────────
    try:
        from PIL import Image
        import pytesseract

        image = Image.open(BytesIO(image_bytes))
        text = pytesseract.image_to_string(image, config="--psm 1")
        if text.strip():
            return [
                Document(
                    page_content=text.strip(),
                    metadata={
                        "source": source_name,
                        "page": "1",
                        "file_type": extension.lstrip("."),
                        "content_type": "ocr_text",
                        "ocr_engine": "tesseract",
                    },
                )
            ]
    except ImportError:
        pass
    except Exception:
        pass

    # ── Gemini Vision fallback ────────────────────────────────────────────────
    api_key = os.getenv("GOOGLE_API_KEY", "")
    if api_key:
        try:
            import base64
            import google.generativeai as genai

            genai.configure(api_key=api_key)
            model = genai.GenerativeModel("gemini-2.0-flash")
            mime_map = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
                        ".tiff": "image/tiff", ".tif": "image/tiff",
                        ".bmp": "image/bmp", ".webp": "image/webp"}
            mime_type = mime_map.get(extension, "image/png")
            img_b64 = base64.b64encode(image_bytes).decode()
            response = model.generate_content([
                "Extract all readable text from this image. Return only the extracted text.",
                {"inline_data": {"mime_type": mime_type, "data": img_b64}},
            ])
            text = response.text.strip()
            if text:
                return [
                    Document(
                        page_content=text,
                        metadata={
                            "source": source_name,
                            "page": "1",
                            "file_type": extension.lstrip("."),
                            "content_type": "ocr_text",
                            "ocr_engine": "gemini",
                        },
                    )
                ]
        except Exception:
            pass

    raise ValueError(
        f"Could not extract text from '{source_name}'. "
        "Install Tesseract + pytesseract, or set GOOGLE_API_KEY for Gemini Vision OCR."
    )


def load_json_document(file_obj, source_name):
    raw_text = _to_bytes(file_obj).decode("utf-8", errors="ignore")
    parsed = json.loads(raw_text)
    pretty = json.dumps(parsed, indent=2, ensure_ascii=True)

    return [
        Document(
            page_content=pretty,
            metadata={
                "source": source_name,
                "page": "JSON",
                "file_type": "json",
                "content_type": "json_document",
            },
        )
    ]


def load_uploaded_file(uploaded_file):
    source_name = getattr(uploaded_file, "name", "uploaded_file")
    extension = Path(source_name).suffix.lower()

    if extension not in SUPPORTED_FILE_TYPES:
        raise ValueError(
            f"Unsupported file type '{extension}'. Supported types: {', '.join(sorted(SUPPORTED_FILE_TYPES))}"
        )

    if extension == ".pdf":
        documents = load_pdf(uploaded_file, source_name)
    elif extension == ".pptx":
        documents = load_pptx(uploaded_file, source_name)
    elif extension in {".xlsx", ".xlsm"}:
        documents = load_excel(uploaded_file, source_name, extension)
    elif extension in {".csv", ".tsv"}:
        documents = load_delimited_text(uploaded_file, source_name, extension)
    elif extension in {".txt", ".md"}:
        documents = load_text_document(uploaded_file, source_name, extension)
    elif extension == ".json":
        documents = load_json_document(uploaded_file, source_name)
    elif extension in IMAGE_EXTENSIONS:
        documents = load_image_ocr(uploaded_file, source_name)
    else:
        documents = []

    if not documents:
        raise ValueError(f"No readable content was found in '{source_name}'.")

    return documents


def extract_tabular_assets(uploaded_file):
    source_name = getattr(uploaded_file, "name", "uploaded_file")
    extension = Path(source_name).suffix.lower()

    if extension not in TABULAR_EXTENSIONS:
        return []

    assets = []
    if extension in {".xlsx", ".xlsm"}:
        frames = _read_excel_frames(uploaded_file)
        for sheet_name, dataframe in frames:
            assets.append(
                {
                    "asset_id": f"{source_name}::{sheet_name}",
                    "source": source_name,
                    "sheet": sheet_name,
                    "file_type": extension.lstrip("."),
                    "dataframe": dataframe,
                    "columns": [str(column) for column in dataframe.columns],
                    "row_count": len(dataframe.index),
                }
            )
    else:
        dataframe = _read_delimited_frame(uploaded_file, extension)
        assets.append(
            {
                "asset_id": f"{source_name}::Data",
                "source": source_name,
                "sheet": "Data",
                "file_type": extension.lstrip("."),
                "dataframe": dataframe,
                "columns": [str(column) for column in dataframe.columns],
                "row_count": len(dataframe.index),
            }
        )

    return assets
